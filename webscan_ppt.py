from pandas import read_excel
import pandas as pd
import pdfplumber as pdfp
import numpy as np
from datetime import datetime
import time
import traceback
import re
from webscan_doc import WebscanParag, WebscanTable
from _G import log_debug, log_error, log_warning, log_info
import _G
import util
import os
import pptx
from pptx import Presentation
from pptx.util import Pt,Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
def generate_ppt_report(save_path, ppt_name, excel):
    global cur_table, cur_parag
    excel = pd.read_excel(excel,sheet_name = 'Developer')
    ppt = Presentation(ppt_name)
    _G.allFileList = os.listdir(_G.pdf_path)
    _G.ppt_params['webCnt']=len(_G.allFileList)
    _G.append_pipe_message(_G.MsgInit)
    for i in _G.allFileList:
        _G.file_list.append(i)
    
    start_t = time.time()
    _G.append_pipe_message(_G.MsgHostList)
    test_target()
    log_info("start_url:", _G.start_url)
    log_debug(f"test_target: {time.time() - start_t}")
 
    start_t = time.time()
    _G.append_pipe_message(_G.MsgRisksCnt)
    risk_cnt()
    log_info("risk:", _G.risk)
    log_debug(f"test_target: {time.time() - start_t}")
 
    start_t = time.time()
    _G.append_pipe_message(_G.MsgHMRisk)
    riskweaklist()
    log_debug(f"riskweaklist: {time.time() - start_t}")
    
    start_t = time.time()
    _G.append_pipe_message(_G.MsgPPTGen)
    ch_text(ppt)
    log_debug(f"ch_text: {time.time() - start_t}")
    
    start_t = time.time()
    chg_table(ppt, excel)
    log_debug(f"chg_table: {time.time() - start_t}")
    

    
    savefile(save_path, ppt)
    
def silent_generate_ppt_report(save_path ,ppt_name, excel):
    try:
        generate_ppt_report(save_path,ppt_name, excel)
    except Exception as err:
        err_info = traceback.format_exc()
        _G.PipeError.append([err, err_info])
        log_error(f"An error occurred during generating report!\n{err_info}")

def move_slide( presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

def delete_slide( presentation,  index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[index])   

def test_target():

    for i in _G.allFileList:
        print(i)
        tmp=[]
        pdf_filename=i
        dev_pdf = pdfp.open(_G.pdf_path+'/'+i)
        dev_pdf = dev_pdf.pages[1]
        test_table = dev_pdf.extract_table()
    #     print(len(test_table))
        for i in range(len(test_table)):
    #         print(i)
            for j in range(5):
    #             print(j)
                if test_table[i][j] != None:
                    if test_table[i][j] != '' :
                        tmp.append(test_table[i][j])
        for i in range(len(tmp)):
            if tmp[i] == 'Start url':
                _G.start_url.append(tmp[i+1])
                
def test_target_table(page, ppt):
        
    # 設置表格位置和大小
    left, top, width, height = Cm(3), Cm(2.5), Cm(3), Cm(3)
    
    second_table_cnt = len(_G.allFileList)

    # 填寫變量內容
    content_arr=[]
    for i in range(0,len(_G.allFileList)):
        content_arr.append([i+1,_G.start_url[i]])
    print(content_arr)

    tmp_page = page
    while second_table_cnt:
        k=1

        if tmp_page != page:
            tmp_page = page
            _G.slide_add+=1
            
            slide = ppt.slides.add_slide(ppt.slide_layouts[1] )

            shape = slide.shapes.add_table(13 if len(_G.allFileList) > 12 else len(_G.allFileList)+1, 2, left, top, width, height)
            table = shape.table
            table.columns[0].width = Cm(2)
            # table.columns[1].width = Cm(2)
            table.columns[1].width = Cm(15)
            table.cell(0, 0).text = "編號"
            # table.cell(0, 1).text = "名稱"
            table.cell(0, 1).text = "URL"
            j=1
            for i in range((_G.range_value*k),(_G.range_value*k)+second_table_cnt):

                for column_index in range(0,len(table.columns)):
                    # 獲取單元格物件
                    cell_temp = table.cell(j, column_index)

    #                 for i in range((range_value*k)+1,(range_value*k)+second_table_cnt):
                    cell_temp.text = str(content_arr[i][column_index])
                j+=1
                # 設定資料

            move_slide(ppt,-1,tmp_page)
        else:

            slide = ppt.slides[page]
            shape = slide.shapes.add_table(13 if len(_G.allFileList) > 12 else len(_G.allFileList)+1, 2, left, top, width, height)
            a = 13 if len(_G.allFileList) > 13 else len(_G.allFileList)
            print('-------')
            print(a)
            table = shape.table
            table.columns[0].width = Cm(2)
            # table.columns[1].width = Cm(2)
            table.columns[1].width = Cm(15)
            table.cell(0, 0).text = "編號"
            # table.cell(0, 1).text = "名稱"
            table.cell(0, 1).text = "URL"
            # print(b)
            for row_index in range(1,len(table.rows)):
                for column_index in range(0,len(table.columns)):
                    # if row_index == 0:
                    #     if column_index == 1:
                    #         cell_temp = table.cell(row_index, column_index)
                    # 獲取單元格物件
                    cell_temp = table.cell(row_index, column_index)

                    # 設定資料
                    cell_temp.text = str(content_arr[row_index-1][column_index])
        second_table_cnt= second_table_cnt - 12 if second_table_cnt > 12 else 0
        page+=1
        k+=1


def risk_cnt():
    for i in _G.allFileList:
        print(i)
        tmp=[]
        pdf_filename=i
        dev_pdf = pdfp.open(_G.pdf_path+'/'+i)
        dev_pdf = dev_pdf.pages[1]
        test_table = dev_pdf.extract_table()
    #     print(len(test_table))
        for i in range(len(test_table)):
    #         print(i)
            for j in range(5):
    #             print(j)
                if test_table[i][j] != None:
                    if test_table[i][j] != '' :
                        tmp.append(test_table[i][j])

        for i in range(len(tmp)):
            if tmp[i] == 'Start url':
                _G.risk.append(tmp[i+1])
            elif tmp[i]=='High':
                _G.high_risk.append(tmp[i+1])
                _G.risk.append(tmp[i+1])
            elif tmp[i]=='Medium':
                _G.mid_risk.append(tmp[i+1])
                _G.risk.append(tmp[i+1])
            elif tmp[i]=='Low':
                _G.risk.append(tmp[i+1])
            elif tmp[i]=='Informational':
                _G.risk.append(tmp[i+1])



def risk_cnt_table(page, ppt):
     # 設置表格位置和大小
    left, top, width, height = Cm(3), Cm(2.5), Cm(3), Cm(3)


    second_table_cnt = len(_G.allFileList)

    # 填寫變量內容
    content_arr=[]
    for i in range(0,len(_G.risk),5):
        content_arr.append([_G.risk[i],_G.risk[i+1],_G.risk[i+2],_G.risk[i+3],_G.risk[i+4]])
    print(content_arr)

    tmp_page = page
    while second_table_cnt:
        k=1

        if tmp_page != page:
            _G.slide_add+=1
            tmp_page = page
            slide = ppt.slides.add_slide(ppt.slide_layouts[1] )

            shape = slide.shapes.add_table(13 if len(_G.allFileList) > 12 else len(_G.allFileList)+1, 5, left, top, width, height)
            # 獲取table對象
            table = shape.table
                 # 設置列寬
            table.columns[0].width = Cm(13)
            table.columns[1].width = Cm(2)
            table.columns[2].width = Cm(2)
            table.columns[3].width = Cm(2)
            table.columns[4].width = Cm(2)
                 # 填寫標題
            table.cell(0, 0).text = "網址"
            table.cell(0, 1).text = "高"
            table.cell(0, 2).text = "中"
            table.cell(0, 3).text = "低"
            table.cell(0, 4).text = "資訊"
            j=1
            for i in range((_G.range_value*k),(_G.range_value*k)+second_table_cnt):

                for column_index in range(0,len(table.columns)):
                    # 獲取單元格物件
                    cell_temp = table.cell(j, column_index)

    #                 for i in range((range_value*k)+1,(range_value*k)+second_table_cnt):
                    cell_temp.text = str(content_arr[i][column_index])
                j+=1
                # 設定資料

            move_slide(ppt,-1,tmp_page)
        else:

            slide = ppt.slides[page]
            shape = slide.shapes.add_table(13 if len(_G.allFileList) > 12 else len(_G.allFileList)+1, 5, left, top, width, height)
            table = shape.table
            # 設置列寬
            table.columns[0].width = Cm(13)
            table.columns[1].width = Cm(2)
            table.columns[2].width = Cm(2)
            table.columns[3].width = Cm(2)
            table.columns[4].width = Cm(2)
            # 填寫標題
            table.cell(0, 0).text = "網址"
            table.cell(0, 1).text = "高"
            table.cell(0, 2).text = "中"
            table.cell(0, 3).text = "低"
            table.cell(0, 4).text = "資訊"
            # print(b)
            for row_index in range(1,len(table.rows)):
                for column_index in range(0,len(table.columns)):
                    # 獲取單元格物件
                    cell_temp = table.cell(row_index, column_index)
                    # 設定資料
                    cell_temp.text = str(content_arr[row_index-1][column_index])
        second_table_cnt= second_table_cnt - 12 if second_table_cnt > 12 else 0
        page+=1
        k+=1
    
#添加表格
def chg_table(ppt, excel):
    for i,slide in enumerate(ppt.slides):
        #     print(i)
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
    #                 print(i)
                    j=0
                    for a  in _G.ppt_table:
                        if  paragraph.text.find(a) !=-1:
                            paragraph.text = paragraph.text.replace(a,'')
                            if a == 'Test_target':
                                test_target_table(i,ppt)
                            elif a == 'testRiskCnt':
                                risk_cnt_table(i+_G.slide_add, ppt)
                            elif a == 'highRiskWeak':
                                highriskweaktable(i+_G.slide_add,ppt)
                            elif a == 'midRiskWeak':
                                midriskweaktable(i+_G.slide_add,ppt)
                            elif a == 'high_risk_explain':
                                _G.append_pipe_message(_G.MsgHighRiskSolve)
                                solvproposal(i+_G.slide_add, ppt, excel )
                                print(a)
                        j+=1
                        

#交換文本  
def ch_text(ppt):
    list_name = list(_G.ppt_params.keys())
    list_target =list(_G.ppt_params.values()) 
    for j in list_target:
        print(j)
    for i,slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    #print(i)
                    i=0
                    for a  in list_name:
                        # print(a)
                        if  paragraph.text.find(a) !=-1:
                            paragraph.text = paragraph.text.replace(a,str(list_target[i]))
                            # print(paragraph.text)
                            # print(i)
                        i+=1 

def riskweaklist():
    for i in _G.allFileList:
        print('檔案 : '+i)
        dev_pdf = pdfp.open(_G.pdf_path+'/'+i)
        DevPdfLen = len(dev_pdf.pages)
        print('頁數 : '+str(DevPdfLen))
        PageList=[]
        for PageNum in range(0,DevPdfLen):
            
    #         print(PageNum)
            page = dev_pdf.pages[PageNum]           # 讀取第一頁
    #         print(page)

            text = page.extract_text()    # 取出文字

            pattern = re.compile('Severity')
            if pattern.findall(text) != []:
    #             print(text)
                PageList.append(PageNum)
    #             print(PageNum+1)
    #     print(PageList)
                        
        tmp=[]
        for j in PageList:
            page = dev_pdf.pages[j]
            PdfTable = page.extract_table()
            
            for j in range(len(PdfTable)):
    #             print(PdfTable[j])
                for k in range(len(PdfTable[j])):
                    if PdfTable[j][k] != None:
                        if PdfTable[j][k] != '' :
                            tmp.append(PdfTable[j][k])


        HighRiskName = set()
        MidRiskName = set()
        for q in range(len(tmp)):
            if 'Severity' in tmp[q] :
                if 'High' in tmp[q+1]:
                    HighRiskName.add(tmp[q-1])
                if 'Medium' in tmp[q+1]:
                    if len(tmp[q-1]) < 4:
                        MidRiskName.add(tmp[q-2])
                    else:
                        MidRiskName.add(tmp[q-1])
        
        for name in HighRiskName:
            if name in _G.HighWeak:
                cnt = _G.HighWeak.get(name)
                cnt+=1
                _G.HighWeak[name]=cnt
            else:
                _G.HighWeak.setdefault(name,1)
            
        for name in MidRiskName:
            if name in _G.MidWeak:
                cnt = _G.MidWeak.get(name)
                cnt+=1
                _G.MidWeak[name]=cnt
            else:
                _G.MidWeak.setdefault(name,1)
                
def highriskweaktable(page, ppt):
    
    left, top, width, height = Cm(3), Cm(2.5), Cm(3), Cm(3)
    slide = ppt.slides[page]
    shape = slide.shapes.add_table(11, 3, left, top, width, height)
    table = shape.table

    table.columns[0].width = Cm(2)
    table.columns[1].width = Cm(13)
    table.columns[2].width = Cm(2)

    table.cell(0, 0).text = '編號'
    table.cell(0, 1).text = '未通過項目'
    table.cell(0, 2).text = '數量'

    content_arr = []
    _G.contentname = [x for x in _G.HighWeak.keys()]
    contentcnt = [x for x in _G.HighWeak.values()]

    for i in range(0, (len(table.rows)-1)):
        if len(_G.contentname)<10:
            _G.contentname.append('無')
            contentcnt.append('無')
        content_arr.append([i+1,_G.contentname[i], contentcnt[i]])
    for rows in range(1,len(table.rows)):
        for col in range(0,len(table.columns)):
                cell_temp = table.cell(rows, col)
                cell_temp.text = str(content_arr[rows-1][col])
    
def midriskweaktable(page, ppt):
    
    left, top, width, height = Cm(3), Cm(2.5), Cm(3), Cm(3)
    slide = ppt.slides[page]
    shape = slide.shapes.add_table(11, 3, left, top, width, height)
    table = shape.table

    table.columns[0].width = Cm(2)
    table.columns[1].width = Cm(13)
    table.columns[2].width = Cm(2)

    table.cell(0, 0).text = '編號'
    table.cell(0, 1).text = '未通過項目'
    table.cell(0, 2).text = '數量'

    content_arr = []
    contentname = [x for x in _G.MidWeak.keys()]
    contentcnt = [x for x in _G.MidWeak.values()]

    for i in range(0, (len(table.rows)-1)):
        if len(contentname)<10:
            contentname.append('無')
            contentcnt.append('無')
        content_arr.append([i+1,contentname[i], contentcnt[i]])
    for rows in range(1,len(table.rows)):
        for col in range(0,len(table.columns)):
                cell_temp = table.cell(rows, col)
                cell_temp.text = str(content_arr[rows-1][col])

def solvproposal(page, ppt, excel ):
        
    left, top, width, height = Cm(3), Cm(2.5), Cm(3), Cm(3)
    
    slide = ppt.slides[page]
    SolvDic = {}
    for i in range(len(excel)):
        if excel['風險名稱'][i] in _G.contentname:
            SolvDic.setdefault(excel['風險名稱'][i],[excel['講中文'][i],excel['修補建議'][i]])
    
    SolvDicLen = len(SolvDic)
    for solvname in SolvDic:
        print('page : '+ str(page))
        shape = slide.shapes.add_table(4, 1, left, top, width, height)
        table = shape.table
        table.columns[0].width = Cm(16)
        table.rows[0].height = Cm(0.5)
        # table.rows[2].height = Cm(0.5)
        # table.rows[4].height = Cm(0.5)
        table.rows[1].height = Cm(0.5)
        table.rows[2].height = Cm(6)
        table.rows[3].height = Cm(1)
        if solvname == '無':
            break
        
        table.cell(0, 0).text = '弱點修補'
        # table.cell(2, 0).text = '弱點說明'
        # table.cell(4, 0).text = '改善建議'
        table.cell(2, 0).text = SolvDic.get(solvname)[0]
        print('table.rows[3].height')
        print(table.rows[3].height)
        print(type(table.rows[3].height))
        table.cell(3, 0).text =  SolvDic.get(solvname)[1]
        print( SolvDic.get(solvname)[1])
        table.cell(1, 0).text = solvname
        print( solvname)
        SolvDicLen-=1
        if SolvDicLen == 0:
            break
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        page+=1
        move_slide(ppt,-1,page)


def savefile(filename,ppt):
    _G.append_pipe_message(_G.MsgPPTGen)
    ppt.save(filename)